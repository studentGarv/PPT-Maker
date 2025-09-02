import os
import hashlib
import requests
from typing import List, Dict, Any, Optional
from pathlib import Path
import tempfile
import re
from dataclasses import dataclass

from pptx import Presentation
import PyPDF2
from urllib.parse import urlparse
import ollama

@dataclass
class DocumentChunk:
    content: str
    source: str
    chunk_id: str
    metadata: Dict[str, Any]

class RAGProcessor:
    """Retrieval-Augmented Generation processor for PPT creation"""
    
    def __init__(self, embed_model: str = "nomic-embed-text"):
        self.embed_model = embed_model
        self.chunks: List[DocumentChunk] = []
        self.embeddings_cache: Dict[str, List[float]] = {}
    
    def process_uploaded_files(self, files: List[str]) -> List[DocumentChunk]:
        """Process uploaded files and extract text content"""
        all_chunks = []
        
        for file_path in files:
            if not os.path.exists(file_path):
                continue
                
            file_ext = Path(file_path).suffix.lower()
            
            try:
                if file_ext == '.pptx':
                    chunks = self._extract_from_pptx(file_path)
                elif file_ext == '.pdf':
                    chunks = self._extract_from_pdf(file_path)
                elif file_ext in ['.txt', '.md']:
                    chunks = self._extract_from_text(file_path)
                else:
                    print(f"Unsupported file type: {file_ext}")
                    continue
                
                all_chunks.extend(chunks)
                print(f"Extracted {len(chunks)} chunks from {Path(file_path).name}")
                
            except Exception as e:
                print(f"Error processing {file_path}: {e}")
        
        self.chunks.extend(all_chunks)
        return all_chunks
    
    def process_urls(self, urls: List[str]) -> List[DocumentChunk]:
        """Process URLs and extract content"""
        all_chunks = []
        
        for url in urls:
            try:
                chunks = self._extract_from_url(url)
                all_chunks.extend(chunks)
                print(f"Extracted {len(chunks)} chunks from {url}")
            except Exception as e:
                print(f"Error processing URL {url}: {e}")
        
        self.chunks.extend(all_chunks)
        return all_chunks
    
    def _extract_from_pptx(self, file_path: str) -> List[DocumentChunk]:
        """Extract text from PowerPoint files"""
        chunks = []
        prs = Presentation(file_path)
        
        for i, slide in enumerate(prs.slides):
            slide_content = []
            slide_title = ""
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    if not slide_title and len(text.split()) <= 10:
                        slide_title = text
                    else:
                        slide_content.append(text)
            
            if slide_content or slide_title:
                content = f"Title: {slide_title}\n" + "\n".join(slide_content)
                chunk_id = hashlib.sha256(f"{file_path}_{i}_{content}".encode()).hexdigest()[:16]
                
                chunks.append(DocumentChunk(
                    content=content,
                    source=f"{Path(file_path).name} - Slide {i+1}",
                    chunk_id=chunk_id,
                    metadata={"type": "pptx", "slide_number": i+1, "file": file_path}
                ))
        
        return chunks
    
    def _extract_from_pdf(self, file_path: str) -> List[DocumentChunk]:
        """Extract text from PDF files"""
        chunks = []
        
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                for i, page in enumerate(pdf_reader.pages):
                    text = page.extract_text().strip()
                    
                    if text:
                        # Split into smaller chunks if too long
                        text_chunks = self._split_text(text, max_length=1000)
                        
                        for j, chunk_text in enumerate(text_chunks):
                            chunk_id = hashlib.sha256(f"{file_path}_{i}_{j}_{chunk_text}".encode()).hexdigest()[:16]
                            
                            chunks.append(DocumentChunk(
                                content=chunk_text,
                                source=f"{Path(file_path).name} - Page {i+1}",
                                chunk_id=chunk_id,
                                metadata={"type": "pdf", "page_number": i+1, "chunk": j+1, "file": file_path}
                            ))
        
        except Exception as e:
            print(f"Error reading PDF {file_path}: {e}")
        
        return chunks
    
    def _extract_from_text(self, file_path: str) -> List[DocumentChunk]:
        """Extract text from text files"""
        chunks = []
        
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read().strip()
                
                if content:
                    text_chunks = self._split_text(content, max_length=1000)
                    
                    for i, chunk_text in enumerate(text_chunks):
                        chunk_id = hashlib.sha256(f"{file_path}_{i}_{chunk_text}".encode()).hexdigest()[:16]
                        
                        chunks.append(DocumentChunk(
                            content=chunk_text,
                            source=f"{Path(file_path).name} - Section {i+1}",
                            chunk_id=chunk_id,
                            metadata={"type": "text", "section": i+1, "file": file_path}
                        ))
        
        except Exception as e:
            print(f"Error reading text file {file_path}: {e}")
        
        return chunks
    
    def _extract_from_url(self, url: str) -> List[DocumentChunk]:
        """Extract content from URLs"""
        chunks = []
        
        try:
            response = requests.get(url, timeout=10)
            response.raise_for_status()
            
            # Simple text extraction (you might want to use BeautifulSoup for better HTML parsing)
            content = response.text
            
            # Remove HTML tags for basic text extraction
            import re
            text_content = re.sub(r'<[^>]+>', ' ', content)
            text_content = re.sub(r'\s+', ' ', text_content).strip()
            
            if text_content:
                text_chunks = self._split_text(text_content, max_length=1000)
                
                for i, chunk_text in enumerate(text_chunks):
                    chunk_id = hashlib.sha256(f"{url}_{i}_{chunk_text}".encode()).hexdigest()[:16]
                    
                    chunks.append(DocumentChunk(
                        content=chunk_text,
                        source=f"{urlparse(url).netloc} - Section {i+1}",
                        chunk_id=chunk_id,
                        metadata={"type": "url", "url": url, "section": i+1}
                    ))
        
        except Exception as e:
            print(f"Error fetching URL {url}: {e}")
        
        return chunks
    
    def _split_text(self, text: str, max_length: int = 1000) -> List[str]:
        """Split text into smaller chunks"""
        if len(text) <= max_length:
            return [text]
        
        chunks = []
        sentences = re.split(r'[.!?]+', text)
        current_chunk = ""
        
        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue
            
            if len(current_chunk) + len(sentence) + 1 <= max_length:
                current_chunk += sentence + ". "
            else:
                if current_chunk:
                    chunks.append(current_chunk.strip())
                current_chunk = sentence + ". "
        
        if current_chunk:
            chunks.append(current_chunk.strip())
        
        return chunks
    
    def _get_embedding(self, text: str) -> List[float]:
        """Get embedding for text using Ollama"""
        cache_key = hashlib.sha256(f"{self.embed_model}_{text}".encode()).hexdigest()
        
        if cache_key in self.embeddings_cache:
            return self.embeddings_cache[cache_key]
        
        try:
            response = ollama.embeddings(model=self.embed_model, prompt=text)
            embedding = response["embedding"]
            self.embeddings_cache[cache_key] = embedding
            return embedding
        except Exception as e:
            print(f"Error getting embedding: {e}")
            return []
    
    def _cosine_similarity(self, a: List[float], b: List[float]) -> float:
        """Calculate cosine similarity between two vectors"""
        if not a or not b or len(a) != len(b):
            return 0.0
        
        dot_product = sum(x * y for x, y in zip(a, b))
        norm_a = sum(x * x for x in a) ** 0.5
        norm_b = sum(y * y for y in b) ** 0.5
        
        if norm_a == 0 or norm_b == 0:
            return 0.0
        
        return dot_product / (norm_a * norm_b)
    
    def retrieve_relevant_chunks(self, query: str, top_k: int = 5) -> List[DocumentChunk]:
        """Retrieve most relevant chunks for a query"""
        if not self.chunks:
            return []
        
        query_embedding = self._get_embedding(query)
        if not query_embedding:
            return self.chunks[:top_k]  # Return first chunks if embedding fails
        
        scored_chunks = []
        
        for chunk in self.chunks:
            chunk_embedding = self._get_embedding(chunk.content)
            if chunk_embedding:
                similarity = self._cosine_similarity(query_embedding, chunk_embedding)
                scored_chunks.append((chunk, similarity))
        
        # Sort by similarity and return top_k
        scored_chunks.sort(key=lambda x: x[1], reverse=True)
        return [chunk for chunk, score in scored_chunks[:top_k]]
    
    def generate_context_prompt(self, original_prompt: str, relevant_chunks: List[DocumentChunk]) -> str:
        """Generate enhanced prompt with context from relevant chunks"""
        if not relevant_chunks:
            return original_prompt
        
        context_sections = []
        for chunk in relevant_chunks:
            context_sections.append(f"From {chunk.source}:\n{chunk.content}")
        
        context_text = "\n\n".join(context_sections)
        
        enhanced_prompt = f"""Create a presentation on: {original_prompt}

Use the following reference materials as context and incorporate relevant information:

{context_text}

Based on the above context and the topic "{original_prompt}", create a comprehensive presentation that:
1. Incorporates relevant information from the reference materials
2. Maintains focus on the main topic
3. Provides actionable insights and takeaways
4. Is well-structured and professional
"""
        
        return enhanced_prompt
    
    def clear_cache(self):
        """Clear all cached data"""
        self.chunks.clear()
        self.embeddings_cache.clear()
        print("RAG cache cleared")

def main():
    """Test the RAG processor"""
    processor = RAGProcessor()
    
    # Test with a sample query
    test_files = ["sample.pptx", "sample.pdf"]  # Add your test files
    test_urls = ["https://example.com"]  # Add test URLs
    
    # Process files
    if any(os.path.exists(f) for f in test_files):
        existing_files = [f for f in test_files if os.path.exists(f)]
        chunks = processor.process_uploaded_files(existing_files)
        print(f"Processed {len(chunks)} chunks from files")
    
    # Test retrieval
    query = "artificial intelligence machine learning"
    relevant_chunks = processor.retrieve_relevant_chunks(query, top_k=3)
    
    print(f"Found {len(relevant_chunks)} relevant chunks for query: '{query}'")
    for chunk in relevant_chunks:
        print(f"- {chunk.source}: {chunk.content[:100]}...")

if __name__ == "__main__":
    main()
