import re, json, math, hashlib
from pathlib import Path
from typing import List, Dict, Any, Optional
from collections import Counter
from dataclasses import dataclass
from pptx import Presentation
import ollama  # Uses installed ollama client

# Import config for default models
try:
    from config import MODEL_NAME, get_improved_filename
    DEFAULT_MODEL = MODEL_NAME
    DEFAULT_EMBED_MODEL = "nomic-embed-text"
except ImportError:
    DEFAULT_MODEL = "gpt-oss:20b"
    DEFAULT_EMBED_MODEL = "nomic-embed-text"

@dataclass
class SlideData:
    index: int
    title: str
    bullets: List[str]
    notes: str
    media: List[str]
    word_count: int

# -------------------------
# PPT LOADING & EXTRACTION
# -------------------------

def load_ppt(path: str) -> Presentation:
    return Presentation(path)

def extract_slides(prs: Presentation) -> List[SlideData]:
    slides = []
    for i, s in enumerate(prs.slides):
        title = ""
        bullets = []
        media = []
        notes = s.notes_slide.notes_text_frame.text if s.has_notes_slide else ""
        for shape in s.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if not title and 0 < len(text.split()) <= 12:
                    title = text
                else:
                    bullets.extend(_split_to_bullets(text))
            # 13 == PICTURE in python-pptx (avoid importing enum to keep light)
            if getattr(shape, "shape_type", None) == 13:
                media.append(f"pic_{i}_{len(media)}")
        wc = sum(len(b.split()) for b in bullets)
        slides.append(SlideData(i, title, bullets, notes, media, wc))
    return slides

# -------------------------
# TEXT PROCESSING HELPERS
# -------------------------

def _split_to_bullets(text: str) -> List[str]:
    raw = re.split(r'[\n•\-–]', text)
    out = []
    for r in raw:
        t = re.sub(r'\s+', ' ', r).strip()
        if not t:
            continue
        if len(t.split()) > 25:  # aggressive split for long paragraphs
            parts = re.split(r'(?<=[.!?])\s+', t)
            out.extend([p.strip() for p in parts if p.strip()])
        else:
            out.append(t)
    return out

# -------------------------
# EMBEDDING-BASED DEDUP
# -------------------------

def _cosine(a: List[float], b: List[float]) -> float:
    if not a or not b or len(a) != len(b):
        return 0.0
    dot = sum(x*y for x, y in zip(a, b))
    na = math.sqrt(sum(x*x for x in a)) or 1e-9
    nb = math.sqrt(sum(y*y for y in b)) or 1e-9
    return dot / (na * nb)

def _get_embedding(text: str, model: str, cache: Dict[str, List[float]]) -> List[float]:
    key = hashlib.sha256((model + "|" + text).encode()).hexdigest()
    if key in cache:
        return cache[key]
    try:
        resp = ollama.embeddings(model=model, prompt=text)
        emb = resp["embedding"]
    except Exception:
        emb = []
    cache[key] = emb
    return emb

def deduplicate_with_embeddings(
    slides: List[SlideData],
    embed_model: str,
    threshold: float = 0.85
) -> List[SlideData]:
    cache: Dict[str, List[float]] = {}
    kept: List[SlideData] = []
    last_emb: Optional[List[float]] = None
    for s in slides:
        text = (s.title + " " + " ".join(s.bullets)).strip()
        emb = _get_embedding(text, embed_model, cache)
        if last_emb is not None and emb:
            sim = _cosine(last_emb, emb)
            if sim >= threshold:
                continue
        kept.append(s)
        if emb:
            last_emb = emb
    return kept

# -------------------------
# LLM OUTLINE GENERATION
# -------------------------

LLM_OUTLINE_SYSTEM = """You are an expert presentation editor. You receive extracted raw bullets from an existing deck. Your task:
1. Remove redundancy
2. Improve clarity & concision
3. Organize into a logical narrative
4. Limit each section to 3-6 sharp bullet points (max ~12 words each)
Return ONLY valid JSON with this structure:
{
  "title": "Improved Main Title",
  "sections": [
     {"title": "Section Title", "points": ["Bullet 1", "Bullet 2"] }
  ],
  "conclusion": {"title": "Conclusion Title", "points": ["Key takeaways", "Call to action"] }
}
Rules:
- No extra commentary
- Bullets: sentence fragments, no trailing punctuation
- Min 2 sections (excluding conclusion) if source has enough content
- Avoid repeating words across consecutive bullets
"""

def _build_source_summary(slides: List[SlideData]) -> str:
    lines = []
    for s in slides:
        lines.append(f"Slide {s.index+1}: {s.title or 'Untitled'}")
        for b in s.bullets[:12]:
            lines.append(f"- {b}")
    return "\n".join(lines[:1200])  # cap

def generate_outline_llm(
    slides: List[SlideData],
    model: str,
    fallback_fn
) -> List[Dict[str, Any]]:
    source = _build_source_summary(slides)
    user_prompt = f"""Source extracted bullets:\n{source}\n\nGenerate improved outline JSON now."""
    try:
        resp = ollama.chat(
            model=model,
            messages=[
                {"role": "system", "content": LLM_OUTLINE_SYSTEM},
                {"role": "user", "content": user_prompt}
            ]
        )
        content = resp["message"]["content"]
        json_match = re.search(r'\{.*\}', content, re.DOTALL)
        if not json_match:
            raise ValueError("No JSON found in model response")
        data = json.loads(json_match.group())
        title = data.get("title") or "Improved Presentation"
        sections = data.get("sections") or []
        conclusion = data.get("conclusion") or {"title": "Conclusion", "points": ["Summary", "Next steps"]}
        outline: List[Dict[str, Any]] = [{"type": "title", "title": title}]
        for sec in sections:
            if not isinstance(sec, dict):
                continue
            pts = [p.strip() for p in sec.get("points", []) if p.strip()]
            outline.append({
                "type": "section",
                "title": sec.get("title", "Section"),
                "points": pts[:6]
            })
        outline.append({
            "type": "conclusion",
            "title": conclusion.get("title", "Conclusion"),
            "points": [p.strip() for p in conclusion.get("points", []) if p.strip()][:6] or ["Thank you"]
        })
        return outline
    except Exception as e:
        print(f"[LLM Outline] Fallback due to error: {e}")
        return fallback_fn(slides)

# -------------------------
# HEURISTIC FALLBACK (OLD OUTLINE)
# -------------------------

def heuristic_outline(slides: List[SlideData]) -> List[Dict[str, Any]]:
    all_text = " ".join([" ".join(s.bullets) for s in slides])
    key_terms = [w for w, _ in Counter(_tokens(all_text)).most_common(12)]
    sections = _segment(slides)
    outline = [{"type": "title", "title": _guess_main_title(slides)}]
    for idx, group in enumerate(sections):
        outline.append({
            "type": "section",
            "title": f"{_title_case(key_terms[idx % len(key_terms)]) if key_terms else 'Section'}",
            "points": _compress_points(group)
        })
    outline.append({"type": "conclusion", "title": "Next Steps", "points": ["Summary", "Recommendations", "Call to Action"]})
    return outline

# (Old outline() removed in favor of generate_outline_llm + heuristic fallback)

# -------------------------
# SUPPORT FUNCTIONS (FALLBACK)
# -------------------------

def _segment(slides: List[SlideData]) -> List[List[SlideData]]:
    seg, cur, wc = [], [], 0
    for s in slides:
        cur.append(s)
        wc += s.word_count
        if wc > 140:
            seg.append(cur)
            cur = []
            wc = 0
    if cur:
        seg.append(cur)
    return seg

def _compress_points(group: List[SlideData]) -> List[str]:
    bullets = []
    for s in group:
        bullets.extend(s.bullets)
    seen = set()
    out = []
    for b in bullets:
        key = hashlib.sha256(b.lower().encode()).hexdigest()[:16]  # Use first 16 chars of SHA-256
        if key in seen:
            continue
        seen.add(key)
        out.append(_shrink(b))
    return out[:6]

def _shrink(text: str) -> str:
    words = text.split()
    if len(words) <= 12:
        return text
    return " ".join(words[:12]) + "..."

def _guess_main_title(slides: List[SlideData]) -> str:
    titles = [s.title for s in slides if s.title]
    if not titles:
        return "Improved Presentation"
    return max(titles, key=lambda t: len([w for w in t.split() if len(w) > 3]))

def _title_case(w: str) -> str:
    return w.capitalize()

def _tokens(s: str):
    return set(re.findall(r'\b[a-z0-9]{3,}\b', s.lower()))

# -------------------------
# PPT REBUILD
# -------------------------

def build_new_ppt(outline: List[Dict[str, Any]], template: Optional[str], out_path: str):
    prs = Presentation(template) if template else Presentation()
    for block in outline:
        slide_layout = 0 if block["type"] == "title" else 1
        slide = prs.slides.add_slide(prs.slide_layouts[slide_layout])
        slide.shapes.title.text = block["title"]
        if block["type"] != "title":
            body = slide.shapes.placeholders[1].text_frame
            first = True
            for p in block.get("points", []):
                if first:
                    body.text = p
                    first = False
                else:
                    body.add_paragraph().text = p
    prs.save(out_path)

# -------------------------
# MAIN IMPROVEMENT ENTRY
# -------------------------

def improve_ppt(
    input_path: str,
    output_path: str,
    template: Optional[str] = None,
    outline_model: str = DEFAULT_MODEL,
    embed_model: str = DEFAULT_EMBED_MODEL,
    dedup_threshold: float = 0.85
):
    """
    Improve an existing PPT:
    1. Extract slides & bullets
    2. Deduplicate using embedding cosine similarity
    3. Generate improved outline via LLM (JSON) with fallback
    4. Build new PPT
    """
    prs = load_ppt(input_path)
    slides = extract_slides(prs)
    slides = deduplicate_with_embeddings(slides, embed_model=embed_model, threshold=dedup_threshold)
    outline = generate_outline_llm(slides, model=outline_model, fallback_fn=heuristic_outline)
    build_new_ppt(outline, template, output_path)

if __name__ == "__main__":
    import argparse
    import os
    
    parser = argparse.ArgumentParser(
        description="Improve existing PowerPoint presentations using AI",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  %(prog)s input.pptx                    # Auto-generates timestamped output
  %(prog)s input.pptx output.pptx        # Specify output filename
  %(prog)s old.pptx improved.pptx --outline-model llama3
  %(prog)s presentation.pptx enhanced.pptx --embed-model all-minilm --threshold 0.9
        """
    )
    
    parser.add_argument(
        "input_path",
        help="Path to the input PowerPoint file"
    )
    
    parser.add_argument(
        "output_path",
        nargs='?',  # Make optional
        help="Path for the improved PowerPoint file (default: auto-generated with timestamp)"
    )
    
    parser.add_argument(
        "--outline-model",
        default=DEFAULT_MODEL,
        help=f"Model for outline generation (default: {DEFAULT_MODEL})"
    )
    
    parser.add_argument(
        "--embed-model",
        default=DEFAULT_EMBED_MODEL,
        help=f"Model for embeddings (default: {DEFAULT_EMBED_MODEL})"
    )
    
    parser.add_argument(
        "--template",
        help="Optional PowerPoint template file"
    )
    
    parser.add_argument(
        "--threshold",
        type=float,
        default=0.85,
        help="Deduplication threshold (0.0-1.0, default: 0.85)"
    )
    
    parser.add_argument(
        "-v", "--verbose",
        action="store_true",
        help="Enable verbose output"
    )
    
    args = parser.parse_args()
    
    # Auto-generate output filename if not provided
    if not args.output_path:
        try:
            from config import get_improved_filename
            args.output_path = get_improved_filename(args.input_path)
        except ImportError:
            # Fallback to simple naming
            from datetime import datetime
            base_name = Path(args.input_path).stem
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            args.output_path = f"{base_name}_improved_{timestamp}.pptx"
    
    # Validate input file exists
    if not os.path.exists(args.input_path):
        print(f"❌ Error: Input file '{args.input_path}' not found")
        exit(1)
    
    if args.verbose:
        print(f"Input file: {args.input_path}")
        print(f"Output file: {args.output_path}")
        print(f"Outline model: {args.outline_model}")
        print(f"Embed model: {args.embed_model}")
        print(f"Threshold: {args.threshold}")
        if args.template:
            print(f"Template: {args.template}")
        print()
    
    print(f"🔧 Improving presentation: {args.input_path}")
    
    try:
        improve_ppt(
            input_path=args.input_path,
            output_path=args.output_path,
            template=args.template,
            outline_model=args.outline_model,
            embed_model=args.embed_model,
            dedup_threshold=args.threshold
        )
        
        print(f"✅ Improved presentation saved to: {args.output_path}")
        
    except Exception as e:
        print(f"❌ Error improving presentation: {e}")
        if args.verbose:
            import traceback
            traceback.print_exc()
        exit(1)