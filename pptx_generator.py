from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from typing import List, Dict
import os


class PPTXGenerator:
    """Class for generating PowerPoint presentations using python-pptx"""
    
    def __init__(self):
        self.presentation = None
    
    def create_presentation(self, title: str) -> None:
        """Create a new presentation"""
        self.presentation = Presentation()
        
        # Set slide dimensions
        self.presentation.slide_width = Inches(10)
        self.presentation.slide_height = Inches(7.5)
    
    def add_title_slide(self, title: str, subtitle: str = "") -> None:
        """Add a title slide to the presentation"""
        slide_layout = self.presentation.slide_layouts[0]  # Title slide layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Format title
        title_frame = title_shape.text_frame
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(44)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
        title_paragraph.alignment = PP_ALIGN.CENTER
        
        # Set subtitle if provided
        if subtitle and slide.placeholders:
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.idx == 1:  # Subtitle placeholder
                    placeholder.text = subtitle
                    subtitle_frame = placeholder.text_frame
                    subtitle_paragraph = subtitle_frame.paragraphs[0]
                    subtitle_paragraph.font.size = Pt(24)
                    subtitle_paragraph.font.color.rgb = RGBColor(89, 89, 89)  # Gray
                    subtitle_paragraph.alignment = PP_ALIGN.CENTER
                    break
    
    def add_content_slide(self, title: str, content: List[str]) -> None:
        """Add a content slide with title and bullet points"""
        slide_layout = self.presentation.slide_layouts[1]  # Title and content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Format title
        title_frame = title_shape.text_frame
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(32)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
        
        # Add content
        content_placeholder = slide.placeholders[1]  # Content placeholder
        text_frame = content_placeholder.text_frame
        text_frame.clear()  # Clear existing content
        
        # Add bullet points
        for i, point in enumerate(content):
            if i == 0:
                paragraph = text_frame.paragraphs[0]
            else:
                paragraph = text_frame.add_paragraph()
            
            paragraph.text = point
            paragraph.level = 0
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray
            paragraph.space_after = Pt(12)
    
    def add_section_slide(self, title: str) -> None:
        """Add a section divider slide"""
        slide_layout = self.presentation.slide_layouts[2]  # Section header layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Format title
        title_frame = title_shape.text_frame
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(40)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White
        title_paragraph.alignment = PP_ALIGN.CENTER
        
        # Add background color to slide
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue
    
    def add_two_column_slide(self, title: str, left_content: List[str], right_content: List[str]) -> None:
        """Add a slide with two columns of content"""
        slide_layout = self.presentation.slide_layouts[3]  # Two content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Format title
        title_frame = title_shape.text_frame
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(32)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(0, 51, 102)
        
        # Left column
        left_placeholder = slide.placeholders[1]
        left_text_frame = left_placeholder.text_frame
        left_text_frame.clear()
        
        for i, point in enumerate(left_content):
            if i == 0:
                paragraph = left_text_frame.paragraphs[0]
            else:
                paragraph = left_text_frame.add_paragraph()
            
            paragraph.text = point
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = RGBColor(51, 51, 51)
        
        # Right column
        right_placeholder = slide.placeholders[2]
        right_text_frame = right_placeholder.text_frame
        right_text_frame.clear()
        
        for i, point in enumerate(right_content):
            if i == 0:
                paragraph = right_text_frame.paragraphs[0]
            else:
                paragraph = right_text_frame.add_paragraph()
            
            paragraph.text = point
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = RGBColor(51, 51, 51)
    
    def add_conclusion_slide(self, title: str = "Thank You", content: List[str] = None) -> None:
        """Add a conclusion/thank you slide"""
        slide_layout = self.presentation.slide_layouts[0]  # Title slide layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Format title
        title_frame = title_shape.text_frame
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(48)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(0, 51, 102)
        title_paragraph.alignment = PP_ALIGN.CENTER
        
        # Add content if provided - optimized for single line
        if content and slide.placeholders:
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.idx == 1:
                    text_frame = placeholder.text_frame
                    text_frame.clear()
                    
                    # For conclusion slides, typically use only the first content item
                    # or combine multiple items into one impactful message
                    if len(content) == 1:
                        # Single line - make it prominent
                        paragraph = text_frame.paragraphs[0]
                        paragraph.text = content[0]
                        paragraph.font.size = Pt(28)
                        paragraph.font.color.rgb = RGBColor(89, 89, 89)
                        paragraph.alignment = PP_ALIGN.CENTER
                        paragraph.space_before = Pt(36)
                    else:
                        # Multiple lines - treat as a combined message
                        combined_text = " • ".join(content)
                        paragraph = text_frame.paragraphs[0]
                        paragraph.text = combined_text
                        paragraph.font.size = Pt(24)
                        paragraph.font.color.rgb = RGBColor(89, 89, 89)
                        paragraph.alignment = PP_ALIGN.CENTER
                        paragraph.space_before = Pt(36)
                    break
    
    def save_presentation(self, filename: str) -> bool:
        """Save the presentation to a file"""
        try:
            if not self.presentation:
                raise ValueError("No presentation created")
            
            # Ensure the filename has .pptx extension
            if not filename.endswith('.pptx'):
                filename += '.pptx'
            
            self.presentation.save(filename)
            return True
        except Exception as e:
            print(f"Error saving presentation: {e}")
            return False
